"""
Script per generare la presentazione PPTX — Stato dei Lavori Aprile 2026
Design: Midnight Executive (navy #1E2761, ice blue #CADCFC)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree
import os

FIGURES = "/Users/danieleuras/Documents/GitHub/miralis-hypergraph-imagined-speech/figures"
OUT = "/Users/danieleuras/Documents/GitHub/miralis-hypergraph-imagined-speech/presentations/stato_lavori_apr2026.pptx"

# ── colori ──────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1E, 0x27, 0x61)
ICE       = RGBColor(0xCA, 0xDC, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE    = RGBColor(0xFF, 0x6B, 0x35)
LIGHT_BG  = RGBColor(0xF4, 0xF7, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
RED_ALERT = RGBColor(0xC0, 0x39, 0x2B)
GREEN_OK  = RGBColor(0x27, 0xAE, 0x60)
GRAY_MID  = RGBColor(0x7F, 0x8C, 0x8D)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank

# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_pt > 0:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             bold=False, italic=False, size=18,
             color=DARK_TEXT, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.italic = italic
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.name  = font_name
    return txb

def add_multiline(slide, lines, x, y, w, h,
                  size=14, color=DARK_TEXT, bullet=False,
                  line_spacing=1.2, font_name="Calibri"):
    """lines = list of (text, bold, size_override_or_None)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, False, None)
        text, bold, sz = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if bullet:
            p.level = 0
        run = p.add_run()
        run.text = ("• " if bullet else "") + text
        run.font.bold = bold
        run.font.size = Pt(sz if sz else size)
        run.font.color.rgb = color
        run.font.name = font_name
    return txb

def add_image(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        # placeholder rect
        r = add_rect(slide, x, y, w, h if h else Inches(3), fill_rgb=ICE)
        add_text(slide, f"[img: {os.path.basename(path)}]",
                 x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), Inches(0.4),
                 size=9, color=GRAY_MID)
        return r
    if h:
        return slide.shapes.add_picture(path, x, y, w, h)
    else:
        return slide.shapes.add_picture(path, x, y, w)

def header_bar(slide, title, subtitle=None):
    """Navy header bar full width"""
    bar_h = Inches(1.0)
    add_rect(slide, 0, 0, W, bar_h, fill_rgb=NAVY)
    add_text(slide, title, Inches(0.35), Inches(0.08), Inches(11), Inches(0.55),
             bold=True, size=26, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.35), Inches(0.62), Inches(10), Inches(0.35),
                 size=13, color=ICE, align=PP_ALIGN.LEFT)

def footer(slide, page_num, total=15):
    add_text(slide, f"Daniele Uras — Politecnico di Milano, DEIB — Aprile 2026",
             Inches(0.3), Inches(7.15), Inches(9), Inches(0.3),
             size=9, color=GRAY_MID)
    add_text(slide, f"{page_num}/{total}",
             Inches(12.5), Inches(7.15), Inches(0.7), Inches(0.3),
             size=9, color=GRAY_MID, align=PP_ALIGN.RIGHT)

def add_note(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def pill_badge(slide, text, x, y, w=Inches(1.8), h=Inches(0.35), bg=NAVY, fg=WHITE, size=11):
    add_rect(slide, x, y, w, h, fill_rgb=bg)
    add_text(slide, text, x + Inches(0.05), y + Inches(0.02), w - Inches(0.1), h,
             bold=True, size=size, color=fg, align=PP_ALIGN.CENTER)

def result_row(slide, cols, x, y, row_h=Inches(0.42), widths=None, bg=WHITE, text_color=DARK_TEXT, size=12):
    if widths is None:
        widths = [Inches(2.5)] * len(cols)
    cur_x = x
    for i, (col, cw) in enumerate(zip(cols, widths)):
        add_rect(slide, cur_x, y, cw, row_h, fill_rgb=bg,
                 line_rgb=ICE, line_pt=0.5)
        add_text(slide, str(col), cur_x + Inches(0.05), y + Inches(0.04),
                 cw - Inches(0.1), row_h - Inches(0.04),
                 size=size, color=text_color, align=PP_ALIGN.CENTER, bold=False)
        cur_x += cw

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

# sfondo a gradiente simulato: due rettangoli
add_rect(s, 0, 0, W, H, fill_rgb=NAVY)
add_rect(s, 0, Inches(4.2), W, Inches(3.3), fill_rgb=RGBColor(0x16, 0x1D, 0x4E))

# striscia decorativa
add_rect(s, 0, Inches(4.0), W, Inches(0.05), fill_rgb=ORANGE)

# titolo
add_text(s, "EEG Imagined Speech Decoding",
         Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.9),
         bold=True, size=38, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(s, "con Graph Neural Networks e Hypergraph",
         Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.7),
         bold=False, size=28, color=ICE, align=PP_ALIGN.CENTER, font_name="Calibri")

# divisore
add_rect(s, Inches(4.5), Inches(2.55), Inches(4.3), Inches(0.04), fill_rgb=ORANGE)

# sottotitolo
add_text(s, "Stato dei Lavori — Aprile 2026",
         Inches(0.8), Inches(2.75), Inches(11.5), Inches(0.5),
         bold=False, size=20, color=ICE, align=PP_ALIGN.CENTER)

# info
add_text(s, "Daniele Uras",
         Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.5),
         bold=True, size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Tesi Magistrale — Politecnico di Milano, DEIB",
         Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.4),
         bold=False, size=15, color=ICE, align=PP_ALIGN.CENTER)
add_text(s, "Supervisore: Francesco Trovò",
         Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4),
         bold=False, size=14, color=RGBColor(0xAD, 0xC4, 0xE8), align=PP_ALIGN.CENTER)

add_note(s, "Presentazione dello stato attuale del progetto di tesi.\nObiettivo: costruire un dizionario neurale semantico che decodifichi il parlato immaginato da segnali EEG usando GNN e Hypergraph Neural Networks.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OBIETTIVO + DATASET
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "Obiettivo e Dataset", "Cosa vogliamo fare e con cosa lavoriamo")
footer(s, 2)

# colonna sinistra — obiettivo
add_text(s, "Obiettivo della Tesi", Inches(0.4), Inches(1.15), Inches(5.8), Inches(0.4),
         bold=True, size=16, color=NAVY)
add_rect(s, Inches(0.4), Inches(1.55), Inches(5.8), Inches(0.04), fill_rgb=ICE)

items_obj = [
    "Decodificare il 'parlato immaginato' da segnali EEG",
    "Mappare pattern cerebrali a categorie semantiche (concretezza: CONCR / AZIONE / STATO / ASTRATTO)",
    "Contributo principale: Hypergraph Neural Networks per relazioni multi-elettrodo",
    "Target da letteratura: Li et al. 2025 (DHSLP) → 78% accuracy subject-specific",
]
for i, item in enumerate(items_obj):
    add_rect(s, Inches(0.4), Inches(1.7 + i * 0.75), Inches(0.06), Inches(0.45), fill_rgb=ORANGE)
    add_text(s, item, Inches(0.55), Inches(1.7 + i * 0.75), Inches(5.6), Inches(0.65),
             size=12.5, color=DARK_TEXT)

# colonna destra — dataset
add_text(s, "Dataset EEG", Inches(6.8), Inches(1.15), Inches(6.0), Inches(0.4),
         bold=True, size=16, color=NAVY)
add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.04), fill_rgb=ICE)

stats = [
    ("70", "Soggetti"),
    ("110", "Parole immaginarie"),
    ("59", "Canali EEG"),
    ("256 Hz", "Frequenza di campionamento"),
    ("~1.5s", "Durata epoca (384 campioni)"),
    ("5", "Sessioni per soggetto"),
    ("~38.000", "Trial totali"),
    ("4", "Cluster semantici (concr4)"),
]
for i, (val, lbl) in enumerate(stats):
    row = i // 2
    col = i % 2
    bx = Inches(6.8 + col * 3.0)
    by = Inches(1.7 + row * 0.9)
    add_rect(s, bx, by, Inches(2.7), Inches(0.75),
             fill_rgb=LIGHT_BG, line_rgb=ICE, line_pt=1)
    add_text(s, val, bx + Inches(0.1), by + Inches(0.03), Inches(2.5), Inches(0.38),
             bold=True, size=20, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, lbl, bx + Inches(0.1), by + Inches(0.4), Inches(2.5), Inches(0.3),
             size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)

# chance level note
add_rect(s, Inches(0.4), Inches(5.55), Inches(12.5), Inches(0.55),
         fill_rgb=RGBColor(0xFF, 0xF3, 0xCD), line_rgb=ORANGE, line_pt=1)
add_text(s, "⚠  Chance level: 25.0% (4 classi bilanciate). Tutte le metriche sono Balanced Accuracy (val_bacc / test_bacc).",
         Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.45),
         size=12, color=RGBColor(0x80, 0x50, 0x00))

add_note(s, "Il task è a 4 classi (schema concr4): parole concr etiche, di azione, di stato, astratte.\nChance level = 25% per classi bilanciate.\nMetrica principale: Balanced Accuracy (val_bacc) perché le classi sono leggermente sbilanciate (ASTRATTO=33%, CONCR=13%).")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — IL PROBLEMA: VARIABILITA' INTER-SOGGETTO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "Il Problema Fondamentale", "La variabilità inter-soggetto domina il segnale semantico")
footer(s, 3)

# grande numero
add_rect(s, Inches(0.4), Inches(1.15), Inches(5.5), Inches(3.0),
         fill_rgb=LIGHT_BG, line_rgb=ICE, line_pt=1)
add_text(s, "ε² = 0.85", Inches(0.5), Inches(1.4), Inches(5.3), Inches(1.2),
         bold=True, size=54, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, "varianza spiegata dal soggetto",
         Inches(0.5), Inches(2.5), Inches(5.3), Inches(0.4),
         size=14, color=GRAY_MID, align=PP_ALIGN.CENTER)
add_rect(s, Inches(1.5), Inches(2.95), Inches(3.3), Inches(0.03), fill_rgb=ORANGE)
add_text(s, "ε² = 0.03  per la parola immaginata",
         Inches(0.5), Inches(3.05), Inches(5.3), Inches(0.4),
         size=14, color=RED_ALERT, align=PP_ALIGN.CENTER, bold=True)
add_text(s, "Il 'segnale' della parola è ~28× più piccolo del 'segnale' dell'identità del soggetto.",
         Inches(0.5), Inches(3.6), Inches(5.3), Inches(0.6),
         size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# implicazioni
add_text(s, "Implicazioni pratiche", Inches(6.3), Inches(1.15), Inches(6.5), Inches(0.4),
         bold=True, size=16, color=NAVY)
add_rect(s, Inches(6.3), Inches(1.55), Inches(6.5), Inches(0.04), fill_rgb=ICE)

impl = [
    ("Cross-subject puro", "Ogni soggetto = dominio separato. Il modello vede 'impronte digitali cerebrali', non parole.", RED_ALERT),
    ("Subject-specific", "Il segnale esiste per alcuni soggetti, ma 330 trial/soggetto sono pochissimi.", ORANGE),
    ("Domain adaptation", "Rimuovere il segnale soggetto (adversarial) porta al collasso totale: non rimane niente.", RED_ALERT),
    ("Meta-learning", "Approccio promettente: impara ad adattarsi velocemente a nuovo soggetto con pochi trial.", GREEN_OK),
]
for i, (title, desc, col) in enumerate(impl):
    by = Inches(1.75 + i * 1.1)
    add_rect(s, Inches(6.3), by, Inches(0.08), Inches(0.85), fill_rgb=col)
    add_text(s, title, Inches(6.5), by, Inches(6.1), Inches(0.38),
             bold=True, size=13, color=col)
    add_text(s, desc, Inches(6.5), by + Inches(0.35), Inches(6.1), Inches(0.5),
             size=11, color=DARK_TEXT)

# ref
add_text(s, "Fonte: Bomatter et al. 2025 — servono >200 soggetti per generalizzare cross-subject su EEG.",
         Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.4),
         size=10, color=GRAY_MID, italic=True)

add_note(s, "Questo è il finding principale del progetto, emerso già nelle analisi esplorative.\nε²=0.85 significa che l'85% della varianza tra trial EEG è spiegata dall'identità del soggetto, non dalla parola immaginata.\nConfronto: in fMRI il rapporto è molto più favorevole alla parola — EEG è intrinsecamente più rumoroso e soggetto-dipendente.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MAPPA ESPERIMENTI
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "Mappa degli Esperimenti", "Progressione dalla baseline alle HGNN")
footer(s, 4)

# timeline orizzontale
notebooks = [
    ("EEG_05\nEEG_07", "Baseline\nBraindecode\n(cross-subj)", "Cross-Subject\n5-Fold CV\n6 modelli", "~25%"),
    ("EEG_06", "Baseline\nSubject-\nSpecific", "Leave-One-\nSession-Out\n6 modelli", "~25-26%"),
    ("EEG_08", "ChebGCN\nGrafo Funz.", "PCC/PLV/wPLI\n5 metodi × 3 arch\nCross-subj", "26.1%"),
    ("EEG_09", "GAT +\nAdversarial", "GRL per\nrimuovere\nsoggetto", "25.4%"),
    ("EEG_10", "Advanced\nGNN", "LGGNet\nDiffPool\nAT-DGNN", "25.9%"),
    ("EEG_11", "Hypergraph\nHGNN", "HGNN_2L\nHGNN_DYN\nHGNN_ATT", "25.0%"),
]

colors_nb = [RGBColor(0x2E, 0x86, 0xC1), RGBColor(0x17, 0xA5, 0x89),
             NAVY, RGBColor(0x8E, 0x44, 0xAD), RED_ALERT, ORANGE]

box_w = Inches(1.9)
box_h = Inches(4.4)
start_x = Inches(0.35)
start_y = Inches(1.2)

for i, (nb, family, details, best) in enumerate(notebooks):
    bx = start_x + i * (box_w + Inches(0.1))
    col = colors_nb[i]

    # header box
    add_rect(s, bx, start_y, box_w, Inches(0.65), fill_rgb=col)
    add_text(s, nb, bx + Inches(0.05), start_y + Inches(0.05),
             box_w - Inches(0.1), Inches(0.55),
             bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # body
    add_rect(s, bx, start_y + Inches(0.65), box_w, box_h - Inches(0.65),
             fill_rgb=LIGHT_BG, line_rgb=col, line_pt=1.5)
    add_text(s, family, bx + Inches(0.08), start_y + Inches(0.75),
             box_w - Inches(0.16), Inches(0.75),
             bold=True, size=12, color=col, align=PP_ALIGN.CENTER)
    add_text(s, details, bx + Inches(0.08), start_y + Inches(1.55),
             box_w - Inches(0.16), Inches(1.0),
             size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # best result
    add_rect(s, bx + Inches(0.2), start_y + Inches(2.75), box_w - Inches(0.4), Inches(0.55),
             fill_rgb=WHITE, line_rgb=col, line_pt=1)
    add_text(s, "Best bacc", bx + Inches(0.2), start_y + Inches(2.8),
             box_w - Inches(0.4), Inches(0.2),
             size=8, color=GRAY_MID, align=PP_ALIGN.CENTER)
    add_text(s, best, bx + Inches(0.2), start_y + Inches(2.98),
             box_w - Inches(0.4), Inches(0.28),
             bold=True, size=14, color=col, align=PP_ALIGN.CENTER)

    # freccia
    if i < len(notebooks) - 1:
        ax = bx + box_w + Inches(0.02)
        ay = start_y + Inches(1.8)
        add_text(s, "→", ax, ay, Inches(0.1), Inches(0.4),
                 bold=True, size=20, color=GRAY_MID, align=PP_ALIGN.CENTER)

# chance line label
add_rect(s, Inches(0.35), Inches(6.0), Inches(12.5), Inches(0.5),
         fill_rgb=RGBColor(0xFF, 0xF0, 0xF0), line_rgb=RED_ALERT, line_pt=0.8)
add_text(s, "⚠  Tutti i modelli convergono a chance (25.0%). Nessuna architettura supera il livello caso cross-subject.",
         Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.4),
         bold=True, size=12, color=RED_ALERT)

add_note(s, "Schema cronologico degli esperimenti.\nEEG_05/07: baseline con architetture standard Braindecode (EEGNet, Deep4Net, EEGConformer, etc.)\nEEG_06: subject-specific — stessa architettura ma train/test sullo stesso soggetto (sessioni diverse)\nEEG_08: primo contributo originale — ChebGCN su grafo funzionale (PCC/PLV/wPLI)\nEEG_09: tentativo domain adaptation con GRL (Gradient Reversal Layer)\nEEG_10: modelli avanzati (LGGNet, AT-DGNN, DiffPool)\nEEG_11: Hypergraph Neural Networks — contributo principale della tesi\nTutti convergono a 25% (chance). Il problema è strutturale.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EEG_05/07 BASELINE CROSS-SUBJECT
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "EEG_05 / EEG_07 — Baseline Cross-Subject", "6 architetture Braindecode, 5-Fold CV, schema concr4")
footer(s, 5)

# tabella risultati
headers = ["Modello", "mean test_bacc", "std", "val_bacc", "Note"]
widths5 = [Inches(2.5), Inches(2.0), Inches(1.5), Inches(2.0), Inches(4.5)]

result_row(s, headers, Inches(0.4), Inches(1.2),
           row_h=Inches(0.45), widths=widths5,
           bg=NAVY, text_color=WHITE, size=12)

rows5 = [
    ("Labram",         "0.2505", "±0.0016", "0.2559", "Tecnicamente above chance (+0.05%): irrilevante"),
    ("ATCNet",         "0.2500", "±0.0000", "0.2507", "Collasso esatto"),
    ("Deep4Net",       "0.2500", "±0.0008", "0.2547", ""),
    ("EEGConformer",   "0.2500", "±0.0000", "0.2500", "Collasso esatto"),
    ("EEGNet",         "0.2500", "±0.0002", "0.2521", ""),
    ("ShallowFBCSPNet","0.2499", "±0.0018", "0.2581", "val_bacc più alta — overfitting"),
]
for i, row in enumerate(rows5):
    bg5 = LIGHT_BG if i % 2 == 0 else WHITE
    result_row(s, row, Inches(0.4), Inches(1.65 + i * 0.44),
               row_h=Inches(0.44), widths=widths5, bg=bg5, size=11)

# chance line annotation
add_rect(s, Inches(0.4), Inches(4.65), Inches(12.5), Inches(0.03), fill_rgb=RED_ALERT)
add_text(s, "Chance = 25.0%", Inches(10.5), Inches(4.55), Inches(2.0), Inches(0.3),
         size=10, color=RED_ALERT, bold=True)

# diagnosi
add_text(s, "Diagnosi", Inches(0.4), Inches(4.85), Inches(12.5), Inches(0.35),
         bold=True, size=14, color=NAVY)

diag_items = [
    "Varianza tra fold ≤ 0.0009 → comportamento stabile e sistematico, non casualità",
    "val_bacc artificialmente alta (0.25–0.26): classi sbilanciate (ASTRATTO=33%) — la balanced accuracy è la metrica corretta",
    "Class collapse: il modello converge a predire la classe maggioritaria per tutti gli input",
    "Conclusione: architetture standard non possono generalizzare cross-soggetto su imagined speech a 4 classi",
]
for i, item in enumerate(diag_items):
    add_rect(s, Inches(0.4), Inches(5.25 + i * 0.42), Inches(0.06), Inches(0.3), fill_rgb=NAVY)
    add_text(s, item, Inches(0.55), Inches(5.25 + i * 0.42), Inches(12.3), Inches(0.38),
             size=11, color=DARK_TEXT)

add_note(s, "EEG_07 è il 5-fold cross-validation di EEG_05. I risultati sono quasi identici tra fold (std ~0.001), il che indica che il fallimento è sistematico e robusto.\nLa differenza tra val_bacc e test_bacc (0.255 vs 0.250) è un artefatto dello sbilanciamento delle classi, non un segnale reale.\nBomatter et al. 2025: servono >200 soggetti per ottenere generalizzazione cross-subject — noi ne abbiamo 70.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — EEG_06 SUBJECT-SPECIFIC
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "EEG_06 — Subject-Specific Baseline", "Leave-One-Session-Out: train sess 1–4, test sess 5. 10 soggetti, concr4.")
footer(s, 6)

# immagine boxplot
img6 = os.path.join(FIGURES, "eeg06_ss_4_norm_boxplot.png")
add_image(s, img6, Inches(0.4), Inches(1.15), Inches(6.5), Inches(4.5))

# analisi destra
add_text(s, "Risultati e Interpretazione", Inches(7.2), Inches(1.15), Inches(5.8), Inches(0.4),
         bold=True, size=15, color=NAVY)
add_rect(s, Inches(7.2), Inches(1.55), Inches(5.8), Inches(0.04), fill_rgb=ICE)

findings6 = [
    ("Mediane ~25%", "Tutti i modelli essenzialmente a chance anche subject-specific.", ORANGE),
    ("ATCNet sotto chance", "Architettura sovradimensionata per 330 trial — convergenza problematica.", RED_ALERT),
    ("Deep4Net outlier 34.4%", "Un soggetto produce bacc=0.344. Il segnale ESISTE per alcuni soggetti, ma non è universale.", GREEN_OK),
    ("Instance Norm applicata", "Normalizzazione per sottrarre impronta soggetto — non basta.", GRAY_MID),
    ("330 trial = troppo pochi", "3 sessioni × ~110 trial. Per 4-class EEG noisy servono >1000 trial/soggetto.", ORANGE),
]
for i, (title, desc, col) in enumerate(findings6):
    by = Inches(1.75 + i * 0.95)
    add_rect(s, Inches(7.2), by, Inches(0.06), Inches(0.72), fill_rgb=col)
    add_text(s, title, Inches(7.35), by, Inches(5.6), Inches(0.35),
             bold=True, size=12, color=col)
    add_text(s, desc, Inches(7.35), by + Inches(0.33), Inches(5.6), Inches(0.45),
             size=10.5, color=DARK_TEXT)

# conclusione box
add_rect(s, Inches(7.2), Inches(6.65), Inches(5.8), Inches(0.6),
         fill_rgb=RGBColor(0xEA, 0xF2, 0xFF), line_rgb=NAVY, line_pt=1.5)
add_text(s, "Il problema non è solo cross-subject: anche subject-specific il segnale imagined speech è debole e inconsistente.",
         Inches(7.35), Inches(6.7), Inches(5.5), Inches(0.5),
         bold=True, size=11, color=NAVY)

add_note(s, "EEG_06 completato il 18 aprile 2026. 10 soggetti, 6 modelli.\nL'outlier Deep4Net a 0.344 è il dato più importante: dimostra che il segnale esiste in alcuni soggetti (il modello 'ci riesce' per quel soggetto), ma non è generalizzabile.\nLa mediana a ~25% suggerisce che la maggior parte dei soggetti non produce un segnale EEG imagined speech decodificabile con questi modelli e questi dati.\nImplicazione per la tesi: il risultato subject-specific basso rafforza la motivazione per il meta-learning (adattamento fast con pochi trial).")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — EEG_07c/07d ANALISI STRUTTURALE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "EEG_07c/07d — Perché Fallisce: Analisi Strutturale", "Prova formale che il grafo PCC non porta informazione semantica cross-subject")
footer(s, 7)

# immagine fingerprint
img7a = os.path.join(FIGURES, "eeg07d_subject_fingerprint.png")
add_image(s, img7a, Inches(0.3), Inches(1.15), Inches(5.8), Inches(3.5))

# immagine laplacian
img7b = os.path.join(FIGURES, "eeg07d_laplacian_spectrum.png")
add_image(s, img7b, Inches(6.3), Inches(1.15), Inches(6.7), Inches(3.5))

# tabella fingerprint
add_text(s, "Subject Fingerprint (cosine similarity)", Inches(0.3), Inches(4.8), Inches(5.8), Inches(0.35),
         bold=True, size=13, color=NAVY)

fp_headers = ["", "Stessa categoria", "Categoria diversa"]
fp_widths  = [Inches(1.8), Inches(1.8), Inches(1.8)]
result_row(s, fp_headers, Inches(0.3), Inches(5.15), widths=fp_widths,
           row_h=Inches(0.38), bg=NAVY, text_color=WHITE, size=10)
result_row(s, ["Stesso soggetto", "0.7851", "0.8014"], Inches(0.3), Inches(5.53),
           widths=fp_widths, row_h=Inches(0.38), bg=LIGHT_BG, size=10)
result_row(s, ["Soggetto diverso", "0.7054", "0.7094"], Inches(0.3), Inches(5.91),
           widths=fp_widths, row_h=Inches(0.38), bg=WHITE, size=10)

add_text(s, "Effetto soggetto: +0.086  |  Effetto categoria: −0.010  →  soggetto domina 8.6×",
         Inches(0.3), Inches(6.35), Inches(5.8), Inches(0.38),
         bold=True, size=11, color=RED_ALERT)

# laplacian annotation
add_text(s, "Spettro Laplaciano per categoria", Inches(6.3), Inches(4.8), Inches(6.7), Inches(0.35),
         bold=True, size=13, color=NAVY)
add_text(s, ("Le 4 curve (CONCR/AZIONE/STATO/ASTRATTO) si sovrappongono perfettamente.\n"
             "Spectral gap identico (~0.15) per tutte le categorie.\n"
             "Implicazione formale: ChebGCN opera nello spazio spettrale → non può fisicamente discriminare categorie.\n"
             "Non è un problema di tuning: è strutturale."),
         Inches(6.3), Inches(5.2), Inches(6.7), Inches(1.5),
         size=11, color=DARK_TEXT)

add_note(s, "EEG_07d è la prova formale del fallimento. Due risultati chiave:\n1) Subject fingerprint: stessa-categoria/soggetti-diversi ha similarity 0.705, stesso-soggetto/categorie-diverse ha 0.801. Il soggetto conta 8.6 volte più della categoria.\n2) Spettro Laplaciano: tutte le categorie hanno la stessa distribuzione degli autovalori del grafo Laplaciano. ChebGCN applica filtri polinomiali sugli autovalori → se gli autovalori sono identici per tutte le categorie, il modello non può distinguerle per definizione.\nQuesto non è un fallimento del tuning o dell'architettura: è un limite matematico.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — EEG_08 ChebGCN RISULTATI
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "EEG_08 — ChebGCN su Grafo Funzionale", "5 metodi grafo × 3 architetture, cross-subject, concr4")
footer(s, 8)

# immagine ablation
img8 = os.path.join(FIGURES, "eeg08_ablation_graph_methods.png")
add_image(s, img8, Inches(0.3), Inches(1.15), Inches(7.5), Inches(4.2))

# tabella risultati
add_text(s, "Best val_bacc per metodo", Inches(8.0), Inches(1.15), Inches(5.0), Inches(0.35),
         bold=True, size=13, color=NAVY)

r8_hdrs = ["Metodo", "Best arch", "val_bacc"]
r8_widths = [Inches(1.5), Inches(2.0), Inches(1.5)]
result_row(s, r8_hdrs, Inches(8.0), Inches(1.55), widths=r8_widths,
           row_h=Inches(0.4), bg=NAVY, text_color=WHITE, size=11)

r8_data = [
    ("PCC",     "ChebGCN_3L",   "0.2607"),
    ("PLV",     "ChebGCN_2L",   "0.2557"),
    ("wPLI",    "ChebGCN_Skip", "0.2535"),
    ("Learned", "ChebGCN_Lea.", "0.2528"),
    ("Dynamic", "ChebGCN_Dyn.", "0.2501"),
]
for i, row in enumerate(r8_data):
    bg8 = LIGHT_BG if i % 2 == 0 else WHITE
    result_row(s, row, Inches(8.0), Inches(1.95 + i * 0.42),
               widths=r8_widths, row_h=Inches(0.42), bg=bg8, size=11)

# osservazioni
add_text(s, "Osservazioni", Inches(8.0), Inches(4.25), Inches(5.0), Inches(0.35),
         bold=True, size=13, color=NAVY)

obs8 = [
    "PCC migliore → ma 26.1% non distinguibile da chance (25%)",
    "Il modello più semplice (2L) spesso migliore → segnale assenza di segnale",
    "Class collapse su STATO: il modello predice sempre la stessa classe",
    "Confusion matrix: colonna STATO piena per tutte le classi reali",
]
for i, o in enumerate(obs8):
    add_rect(s, Inches(8.0), Inches(4.65 + i * 0.44), Inches(0.06), Inches(0.35), fill_rgb=ORANGE)
    add_text(s, o, Inches(8.12), Inches(4.65 + i * 0.44), Inches(4.85), Inches(0.4),
             size=10.5, color=DARK_TEXT)

# architettura nota
add_rect(s, Inches(0.3), Inches(5.55), Inches(7.5), Inches(0.6),
         fill_rgb=LIGHT_BG, line_rgb=ICE, line_pt=0.8)
add_text(s, "Architettura: ispirata a GCNs-Net (IEEE TNSRE 2022) — 1D CNN temporal encoder + ChebConv(K=2) + global_mean_pool + MLP",
         Inches(0.45), Inches(5.6), Inches(7.2), Inches(0.5), size=10.5, color=DARK_TEXT)

add_note(s, "EEG_08 è il primo contributo originale GNN della tesi.\nGrafo PCC: Pearson Correlation Coefficient tra i 59 canali EEG, k-NN con k=6.\nChebConv(K=2): filtri di Chebyshev di ordine 2 — implementazione ispirata a GCNs-Net (Lun et al., IEEE TNSRE 2022).\nIl problema del class collapse: con classi sbilanciate, senza class weights il modello predice sempre ASTRATTO (classe più frequente). Con class weights, predice sempre STATO. In nessun caso impara pattern semantici.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EEG_09 GAT + ADVERSARIAL
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "EEG_09 — GAT + Domain Adversarial Training", "Tentativo di rimuovere il segnale soggetto con Gradient Reversal Layer")
footer(s, 9)

# immagine confusion matrix adversarial
img9 = os.path.join(FIGURES, "eeg09_confmat_GAT_2L_ADV.png")
add_image(s, img9, Inches(0.3), Inches(1.15), Inches(5.5), Inches(3.8))

# immagine risultati
img9b = os.path.join(FIGURES, "eeg09_gat_results_4cls.png")
add_image(s, img9b, Inches(6.0), Inches(1.15), Inches(7.1), Inches(3.8))

# smoking gun
add_rect(s, Inches(0.3), Inches(5.1), Inches(12.8), Inches(0.9),
         fill_rgb=RGBColor(0xFF, 0xEB, 0xEB), line_rgb=RED_ALERT, line_pt=1.5)
add_text(s, "SMOKING GUN", Inches(0.5), Inches(5.15), Inches(2.0), Inches(0.45),
         bold=True, size=14, color=RED_ALERT)
add_text(s, ("L'adversarial training (GAT_2L_ADV) ha rimosso il segnale del soggetto → collasso totale su classe 1 (0.81–0.82 per ogni classe reale). "
             "L'unico segnale discriminativo nei grafi EEG è l'identità del soggetto. "
             "Senza di essa, il modello non ha nulla da imparare."),
         Inches(2.6), Inches(5.15), Inches(10.3), Inches(0.8),
         bold=False, size=11.5, color=RED_ALERT)

# tabella
add_text(s, "Risultati (v2: λ=0.5, patience=40)", Inches(0.3), Inches(6.1), Inches(5.0), Inches(0.35),
         bold=True, size=12, color=NAVY)
r9_hdrs = ["Modello", "val_bacc", "test_bacc"]
r9_widths = [Inches(2.0), Inches(1.5), Inches(1.5)]
result_row(s, r9_hdrs, Inches(0.3), Inches(6.45),
           widths=r9_widths, row_h=Inches(0.38), bg=NAVY, text_color=WHITE, size=11)
r9_data = [("GAT_2L", "0.262", "0.250"),
           ("GAT_2L_ADV", "0.258", "0.254"),
           ("GAT_3L_ADV", "0.258", "0.250")]
for i, row in enumerate(r9_data):
    result_row(s, row, Inches(0.3), Inches(6.83 + i * 0.38),
               widths=r9_widths, row_h=Inches(0.38),
               bg=LIGHT_BG if i%2==0 else WHITE, size=11)

add_note(s, "EEG_09 implementa il domain adversarial training con GRL (Gradient Reversal Layer).\nL'idea: aggiungere un discriminatore che predice l'identità del soggetto dall'embedding. Invertendo il gradiente del discriminatore, forziamo l'encoder a rimuovere l'informazione soggetto.\nRisultato: funziona tecnicamente (alpha GRL raggiunge 1.0) ma porta al collasso totale. Questo conferma che il segnale soggetto era l'UNICA cosa che il modello stava imparando.\nConfronto SOTA: FAST 2025 (Jiang et al., 57 soggetti, 5 classi, strict LOSO): test_bacc=26.9% (+6.9% su chance=20%). Il problema è aperto in letteratura.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — EEG_10 ADVANCED GNN
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "EEG_10 — Advanced GNN", "LGGNet, AT-DGNN, DiffPool — tutti convergono a ln(4) = 1.386")
footer(s, 10)

# immagine comparison
img10 = os.path.join(FIGURES, "eeg10_comparison.png")
add_image(s, img10, Inches(0.3), Inches(1.15), Inches(7.0), Inches(3.8))

# immagine lggnet training
img10b = os.path.join(FIGURES, "eeg10_lggnet_training.png")
add_image(s, img10b, Inches(7.5), Inches(1.15), Inches(5.5), Inches(3.8))

# tabella
add_text(s, "Risultati (concr4, cross-subject)", Inches(0.3), Inches(5.1), Inches(6.5), Inches(0.35),
         bold=True, size=13, color=NAVY)
r10_hdrs = ["Modello", "val_bacc", "Note"]
r10_widths = [Inches(2.2), Inches(1.5), Inches(5.0)]
result_row(s, r10_hdrs, Inches(0.3), Inches(5.5),
           widths=r10_widths, row_h=Inches(0.4), bg=NAVY, text_color=WHITE, size=11)
r10_data = [
    ("ChebGCN (EEG_08)", "0.261", "Baseline — modello più semplice = migliore"),
    ("LGGNet",           "0.259", "Loss → ln(4)=1.386 = output uniforme ottimale"),
    ("AT-DGNN",          "0.254", "Attenzione dinamica → nessun vantaggio"),
    ("DiffPool-GNN",     "0.253", "Pooling gerarchico → nessun vantaggio"),
]
for i, row in enumerate(r10_data):
    result_row(s, row, Inches(0.3), Inches(5.9 + i * 0.4),
               widths=r10_widths, row_h=Inches(0.4),
               bg=LIGHT_BG if i%2==0 else WHITE, size=11)

# insight ln(4)
add_rect(s, Inches(7.5), Inches(5.1), Inches(5.6), Inches(1.9),
         fill_rgb=RGBColor(0xFF, 0xF3, 0xCD), line_rgb=ORANGE, line_pt=1.5)
add_text(s, "Loss = ln(4) = 1.386", Inches(7.65), Inches(5.2), Inches(5.3), Inches(0.45),
         bold=True, size=16, color=ORANGE)
add_text(s,
         ("Il minimo teorico di CrossEntropy per output uniforme su 4 classi è ln(4).\n"
          "LGGNet converge esattamente a questo valore.\n"
          "Il modello ha trovato la soluzione ottima per un task senza segnale:\n"
          "predire con probabilità uniforme 25% per ogni classe."),
         Inches(7.65), Inches(5.65), Inches(5.3), Inches(1.25),
         size=11, color=DARK_TEXT)

add_note(s, "EEG_10 testa 3 architetture avanzate:\n- LGGNet (Lei et al.): graph neural network con lateral graph convolution per EEG motoria. Adattata a imagined speech.\n- AT-DGNN: attention-based dynamic graph neural network.\n- DiffPool: pooling gerarchico differenziabile.\nIl risultato più illuminante è LGGNet: la loss converge esattamente a ln(4) = 1.386, che è il minimo teorico per output uniforme. Il modello non è 'confuso' — sta trovando la soluzione ottimale per un problema senza segnale.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — EEG_11 HYPERGRAPH
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "EEG_11 — Hypergraph Neural Networks", "Contributo principale della tesi: relazioni multi-elettrodo con iperspigoli")
footer(s, 11)

# immagine hyperedge topomaps
img11a = os.path.join(FIGURES, "eeg10_hyperedge_topomaps.png")
add_image(s, img11a, Inches(0.3), Inches(1.15), Inches(5.5), Inches(3.2))

# immagine graph vs hypergraph
img11b = os.path.join(FIGURES, "eeg10_graph_vs_hypergraph.png")
add_image(s, img11b, Inches(6.0), Inches(1.15), Inches(7.1), Inches(3.2))

# descrizione modelli
add_text(s, "3 Architetture HGNN", Inches(0.3), Inches(4.5), Inches(5.5), Inches(0.35),
         bold=True, size=13, color=NAVY)

models_11 = [
    ("HGNN_2L", "Iperspigoli statici (Feng et al. 2019)", "HGNN classico baseline"),
    ("HGNN_2L_DYN", "Layer 2 dinamico (ispirazione Li et al. 2025)", "Iperspigoli ricalcolati per trial"),
    ("HGNN_ATT_2L", "Attention sugli iperspigoli (AllSet, Chien et al.)", "Aggregazione pesata per importanza"),
]
for i, (name, arch, desc) in enumerate(models_11):
    bx = Inches(0.3 + i * 4.3)
    add_rect(s, bx, Inches(4.9), Inches(4.0), Inches(0.4),
             fill_rgb=NAVY if i==0 else (ORANGE if i==1 else GREEN_OK))
    add_text(s, name, bx + Inches(0.05), Inches(4.92), Inches(3.9), Inches(0.36),
             bold=True, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, arch, bx + Inches(0.05), Inches(5.35), Inches(3.9), Inches(0.35),
             size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    add_text(s, desc, bx + Inches(0.05), Inches(5.68), Inches(3.9), Inches(0.35),
             size=9.5, color=GRAY_MID, italic=True, align=PP_ALIGN.CENTER)

# risultati + performance fix
add_text(s, "Risultati", Inches(0.3), Inches(6.15), Inches(5.5), Inches(0.35),
         bold=True, size=12, color=NAVY)
r11_hdrs = ["Modello", "test_bacc", "Pattern"]
r11_widths = [Inches(2.0), Inches(1.3), Inches(2.5)]
result_row(s, r11_hdrs, Inches(0.3), Inches(6.5),
           widths=r11_widths, row_h=Inches(0.35), bg=NAVY, text_color=WHITE, size=10)
r11_data = [("HGNN_2L", "0.250", "Split classi 1+2"),
            ("HGNN_DYN", "0.247", "Collasso su classe 3"),
            ("HGNN_ATT", "0.250", "Collasso su classe 1")]
for i, row in enumerate(r11_data):
    result_row(s, row, Inches(0.3), Inches(6.85 + i * 0.35),
               widths=r11_widths, row_h=Inches(0.35),
               bg=LIGHT_BG if i%2==0 else WHITE, size=10)

# performance fix
add_text(s, "Fix Performance Implementati", Inches(6.0), Inches(6.15), Inches(7.1), Inches(0.35),
         bold=True, size=12, color=NAVY)
pfixes = [
    "Fix 1: hyperedge_index globale (no PCC per-trial) → −7ms/trial",
    "Fix 2: preload RAM (~3.4 GB) → zero I/O H5 durante training",
    "Fix 3: dataset pre-built su disco (~3.7 GB) → training istantaneo a sessioni successive",
]
for i, f in enumerate(pfixes):
    add_rect(s, Inches(6.0), Inches(6.52 + i * 0.38), Inches(0.06), Inches(0.3), fill_rgb=GREEN_OK)
    add_text(s, f, Inches(6.12), Inches(6.52 + i * 0.38), Inches(6.9), Inches(0.35),
             size=11, color=DARK_TEXT)

add_note(s, "EEG_11 è il contributo principale della tesi.\nHypergraph: invece di edge binari (elettrodo A ↔ elettrodo B), gli iperspigoli connettono K elettrodi contemporaneamente (cluster topologico). Questo cattura relazioni multi-nodo che il grafo standard non può rappresentare.\nI 3 modelli implementati:\n- HGNN_2L: baseline con iperspigoli statici calcolati da PCC (Feng et al. 2019)\n- HGNN_2L_DYN: il secondo layer ricalcola gli iperspigoli per ogni trial (ispirazione Li et al. 2025)\n- HGNN_ATT_2L: attention sugli iperspigoli (AllSet, Chien et al. 2022) — pesatura adattiva\nAnche qui, tutti a chance. La motivazione strutturale di EEG_07d spiega perché: gli iperspigoli sono costruiti dallo stesso PCC che ha spettro Laplaciano identico tra categorie.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — QUADRO COMPLESSIVO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "Quadro Complessivo", "Tutti gli esperimenti: cross-subject e subject-specific")
footer(s, 12)

# tabella grande
all_hdrs = ["Notebook", "Famiglia", "Modelli", "Setting", "Best bacc", "Vs Chance"]
all_widths = [Inches(1.3), Inches(2.0), Inches(2.5), Inches(2.2), Inches(1.5), Inches(3.5)]
result_row(s, all_hdrs, Inches(0.3), Inches(1.2),
           widths=all_widths, row_h=Inches(0.45), bg=NAVY, text_color=WHITE, size=11)

all_data = [
    ("EEG_05/07", "Braindecode",    "EEGNet, Deep4Net, EEGConformer, Labram...", "Cross-subj 5-fold", "25.1%", "≈ Chance"),
    ("EEG_06",    "Braindecode",    "EEGNet, Deep4Net, EEGConformer, ATCNet...", "Subject-specific",  "~26%",  "≈ Chance (1 outlier 34.4%)"),
    ("EEG_08",    "ChebGCN",        "2L, 3L, Skip × 5 metodi grafo",            "Cross-subj",        "26.1%", "+1.1% (non sign.)"),
    ("EEG_09",    "GAT + Adversar.","GAT_2L, GAT_2L_ADV, GAT_3L_ADV",           "Cross-subj",        "25.4%", "≈ Chance"),
    ("EEG_10",    "Advanced GNN",   "LGGNet, AT-DGNN, DiffPool",                "Cross-subj",        "25.9%", "≈ Chance (loss=ln4)"),
    ("EEG_11",    "HGNN",           "HGNN_2L, HGNN_DYN, HGNN_ATT",             "Cross-subj",        "25.0%", "= Chance esatto"),
]
for i, row in enumerate(all_data):
    bg_all = LIGHT_BG if i % 2 == 0 else WHITE
    # colora la colonna "vs chance" in rosso/verde
    result_row(s, row, Inches(0.3), Inches(1.65 + i * 0.52),
               widths=all_widths, row_h=Inches(0.52), bg=bg_all, size=10.5)

# confronto SOTA
add_text(s, "Confronto con SOTA", Inches(0.3), Inches(5.0), Inches(12.5), Inches(0.38),
         bold=True, size=14, color=NAVY)
add_rect(s, Inches(0.3), Inches(5.38), Inches(12.5), Inches(0.04), fill_rgb=ICE)

sota_items = [
    "Li et al. 2025 (DHSLP) — TARGET: 78% su imagined speech subject-specific (dataset diverso, setting controllato)",
    "FAST 2025 (Jiang et al.) — 57 soggetti, 5 classi, strict LOSO: 26.9% (+6.9% su chance 20%): il problema è aperto",
    "Bomatter et al. 2025 — cross-subject EEG richiede >200 soggetti per generalizzare: noi ne abbiamo 70",
    "Einizade et al. 2022 (GraphIS) — punto di partenza GCN per imagined speech: risultati in setting meno strict",
]
for i, item in enumerate(sota_items):
    add_rect(s, Inches(0.3), Inches(5.5 + i * 0.44), Inches(0.06), Inches(0.35), fill_rgb=ICE)
    add_text(s, item, Inches(0.45), Inches(5.5 + i * 0.44), Inches(12.3), Inches(0.4),
             size=11, color=DARK_TEXT)

add_note(s, "Tabella riassuntiva di tutti gli esperimenti.\nIl pattern è consistente: nessuna architettura, nessuna strategia di dominio, nessuna topologia grafo riesce a superare chance cross-subject.\nL'unico dato positivo è l'outlier Deep4Net in EEG_06 (subject-specific, 34.4%) — dimostra che il segnale esiste in alcuni soggetti.\nConfronto SOTA: il problema è difficile anche in letteratura. Li et al. 2025 ottengono risultati eccellenti ma in setting molto più controllati (pochi soggetti, molti trial per soggetto).")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DIAGNOSI FINALE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "Diagnosi Finale", "4 prove convergenti del fallimento cross-subject")
footer(s, 13)

diagnosi = [
    ("1", "EEG_07c — Analisi esplorativa",
     "F-score ANOVA degli edge vicino a 1 per tutte le categorie. t-SNE mostra cluster per SOGGETTO, non per parola. Intra ≈ inter distanza per categoria.",
     RED_ALERT),
    ("2", "EEG_07d — Subject fingerprint",
     "Cosine similarity: soggetto domina categoria 8.6×. Spettro Laplaciano identico per tutte le categorie → ChebGCN/HGNN non possono discriminare per definizione matematica.",
     RED_ALERT),
    ("3", "EEG_09 — Smoking gun adversarial",
     "Rimuovendo il segnale del soggetto con GRL, il modello collassa a 0.25. L'unico segnale discriminativo era l'identità del soggetto, non la parola immaginata.",
     RED_ALERT),
    ("4", "EEG_08/10/11 — Convergenza a ln(4)",
     "Tutte le architetture (ChebGCN, LGGNet, AT-DGNN, DiffPool, HGNN) convergono a loss=ln(4)=1.386. Il modello trova la soluzione ottima per un task senza segnale: output uniforme.",
     RED_ALERT),
]

for i, (num, title, desc, col) in enumerate(diagnosi):
    by = Inches(1.25 + i * 1.4)
    add_rect(s, Inches(0.3), by, Inches(0.7), Inches(1.1), fill_rgb=col)
    add_text(s, num, Inches(0.3), by + Inches(0.25), Inches(0.7), Inches(0.6),
             bold=True, size=32, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(1.05), by, Inches(11.8), Inches(1.1),
             fill_rgb=LIGHT_BG, line_rgb=col, line_pt=1)
    add_text(s, title, Inches(1.15), by + Inches(0.08), Inches(11.5), Inches(0.38),
             bold=True, size=14, color=col)
    add_text(s, desc, Inches(1.15), by + Inches(0.46), Inches(11.5), Inches(0.55),
             size=11.5, color=DARK_TEXT)

# conclusione box
add_rect(s, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.55),
         fill_rgb=NAVY)
add_text(s,
         "Conclusione: il fallimento cross-subject è sistematico, strutturale e confermato da 4 assi indipendenti. Non è risolvibile con tuning.",
         Inches(0.45), Inches(6.8), Inches(12.5), Inches(0.45),
         bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)

add_note(s, "Questa slide riassume le 4 prove convergenti del fallimento.\nLa chiave è che le prove vengono da QUATTRO angolazioni diverse e indipendenti:\n1) Analisi statistica pre-training (ANOVA degli edge)\n2) Analisi strutturale del grafo (fingerprint, spettro Laplaciano)\n3) Esperimento di perturbazione (domain adversarial)\n4) Osservazione del comportamento di training (convergenza a ln(4))\nQuando 4 angolazioni diverse convergono alla stessa conclusione, il risultato è robusto.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — META-LEARNING: NEXT STEP
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "Prossimo Passo: Meta-Learning", "Learning to learn — adattamento rapido a nuovo soggetto con pochi trial")
footer(s, 14)

# problema-soluzione
add_text(s, "Perché Meta-Learning?", Inches(0.4), Inches(1.15), Inches(5.8), Inches(0.38),
         bold=True, size=16, color=NAVY)
add_rect(s, Inches(0.4), Inches(1.53), Inches(5.8), Inches(0.04), fill_rgb=ICE)

gap_items = [
    ("Cross-subject puro",   "Fallisce: ε²=0.85, ogni soggetto = dominio separato"),
    ("Subject-specific",     "Fallisce: 330 trial sono pochi per adattarsi"),
    ("Domain adversarial",   "Fallisce: rimuovere soggetto = rimuovere tutto"),
    ("Meta-learning",        "Soluzione: impara a adattarsi con K trial per soggetto"),
]
for i, (approach, result) in enumerate(gap_items):
    col = GREEN_OK if i == 3 else (RED_ALERT if i < 3 else GRAY_MID)
    bg_g = RGBColor(0xEA, 0xF7, 0xEA) if i == 3 else RGBColor(0xFF, 0xEB, 0xEB)
    add_rect(s, Inches(0.4), Inches(1.65 + i * 0.85), Inches(5.8), Inches(0.75),
             fill_rgb=bg_g, line_rgb=col, line_pt=1.5)
    add_text(s, ("✓ " if i == 3 else "✗ ") + approach,
             Inches(0.55), Inches(1.7 + i * 0.85), Inches(5.5), Inches(0.3),
             bold=True, size=12, color=col)
    add_text(s, result,
             Inches(0.55), Inches(2.0 + i * 0.85), Inches(5.5), Inches(0.38),
             size=11, color=DARK_TEXT)

# spiegazione meta-learning
add_text(s, "Come funziona (Prototypical Networks)", Inches(6.6), Inches(1.15), Inches(6.4), Inches(0.38),
         bold=True, size=16, color=NAVY)
add_rect(s, Inches(6.6), Inches(1.53), Inches(6.4), Inches(0.04), fill_rgb=ICE)

ml_steps = [
    ("Meta-train", "60 soggetti. Per ogni episodio: scegli soggetto i, K support trial per classe, Q query trial. Impara embedding che classifica per distanza dal prototipo."),
    ("Meta-test", "10 soggetti MAI visti. Dai K=5 trial per classe → calcola prototipi → classifica query. Nessun gradient step aggiuntivo."),
    ("Vantaggio", "Sfrutta i 70 soggetti come 'compiti diversi'. Impara struttura comune tra soggetti, non memorizza identità."),
]
for i, (phase, desc) in enumerate(ml_steps):
    by = Inches(1.7 + i * 1.4)
    col = [NAVY, ORANGE, GREEN_OK][i]
    add_rect(s, Inches(6.6), by, Inches(1.5), Inches(1.1), fill_rgb=col)
    add_text(s, phase, Inches(6.62), by + Inches(0.3), Inches(1.46), Inches(0.5),
             bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(8.1), by, Inches(4.9), Inches(1.1),
             fill_rgb=LIGHT_BG, line_rgb=col, line_pt=1)
    add_text(s, desc, Inches(8.2), by + Inches(0.12), Inches(4.7), Inches(0.85),
             size=11, color=DARK_TEXT)

# timeline implementazione
add_rect(s, Inches(0.4), Inches(5.85), Inches(12.7), Inches(0.55),
         fill_rgb=LIGHT_BG, line_rgb=ICE, line_pt=1)
add_text(s, "Piano: EEG_12 Meta-Learning — EEGNet encoder + Prototypical Networks, split 60/10 soggetti, episodic dataloader. Stima: 1-2 settimane.",
         Inches(0.55), Inches(5.9), Inches(12.4), Inches(0.45),
         size=12, color=DARK_TEXT)

# narrative tesi
add_text(s, "Narrativa tesi aggiornata:",
         Inches(0.4), Inches(6.5), Inches(12.7), Inches(0.32),
         bold=True, size=12, color=NAVY)
add_text(s,
         "EEG_05-11 dimostrano che né cross-subject né subject-specific funzionano → bottleneck strutturale (ε²=0.85) + dati insufficienti. "
         "Meta-learning è la risposta metodologica: adattamento con K trial per soggetto, sfruttando tutti i 70 soggetti come episodi di training.",
         Inches(0.4), Inches(6.82), Inches(12.7), Inches(0.55),
         size=11, color=DARK_TEXT)

add_note(s, "Il meta-learning è il naturale passo successivo dato il fallimento cross-subject e subject-specific.\nPrototypical Networks (Snell et al. 2017) sono particolarmente adatti per EEG:\n- Non richiedono gradient steps al test time (al contrario di MAML)\n- Funzionano bene con pochi shot\n- L'idea del 'prototipo per classe' è naturale per EEG: prototipo = pattern medio di una categoria per quel soggetto\nAlternativamente: MAML, Reptile, o anche transfer learning fine-tuning da backbone pre-trained.\nBiblioteca: learn2learn (pip install learn2learn) o implementazione custom con episodic DataLoader.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ROADMAP E CONCLUSIONI
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
header_bar(s, "Roadmap e Conclusioni", "Stato attuale, prossimi passi, domanda di ricerca aperta")
footer(s, 15)

# colonna sinistra: completato
add_text(s, "Completato ✓", Inches(0.4), Inches(1.15), Inches(5.8), Inches(0.38),
         bold=True, size=16, color=GREEN_OK)
add_rect(s, Inches(0.4), Inches(1.53), Inches(5.8), Inches(0.04), fill_rgb=ICE)

done_items = [
    "Pipeline dati completa (59 ch, 70 soggetti, ~38k trial)",
    "Baseline cross-subject: 6 modelli Braindecode (EEG_05/07)",
    "Baseline subject-specific con Leave-One-Session-Out (EEG_06)",
    "Analisi esplorativa grafo: ANOVA, t-SNE, fingerprint, Laplaciano (EEG_07c/07d)",
    "ChebGCN su 5 metodi grafo (PCC/PLV/wPLI/learned/dynamic) (EEG_08)",
    "GAT + Domain Adversarial Training (EEG_09)",
    "LGGNet, AT-DGNN, DiffPool-GNN (EEG_10)",
    "Hypergraph Neural Networks: HGNN_2L, DYN, ATT (EEG_11)",
    "Diagnosi strutturale: 4 prove convergenti del fallimento cross-subject",
]
for i, item in enumerate(done_items):
    add_rect(s, Inches(0.4), Inches(1.65 + i * 0.47), Inches(0.06), Inches(0.35), fill_rgb=GREEN_OK)
    add_text(s, item, Inches(0.55), Inches(1.65 + i * 0.47), Inches(5.6), Inches(0.42),
             size=10.5, color=DARK_TEXT)

# colonna destra: prossimi passi
add_text(s, "Prossimi Passi", Inches(6.6), Inches(1.15), Inches(6.4), Inches(0.38),
         bold=True, size=16, color=NAVY)
add_rect(s, Inches(6.6), Inches(1.53), Inches(6.4), Inches(0.04), fill_rgb=ICE)

next_items = [
    ("IMMEDIATO", "EEG_12 — Meta-Learning (Prototypical Networks)\nEncoder EEGNet + few-shot su 10 soggetti test", ORANGE),
    ("BREVE", "Analisi completa EEG_06 su tutti i 70 soggetti\n(attualmente su 10 — identificare soggetti 'buoni')", NAVY),
    ("BREVE", "Ablation study: HGNN vs GNN su subject-specific\n(replicare EEG_08/11 in setting subject-specific)", NAVY),
    ("SCRITTURA", "Capitoli tesi: Metodologia HGNN, Esperimenti,\nDiagnosi strutturale, Meta-learning come futuro", RGBColor(0x6C, 0x35, 0x7B)),
]
for i, (priority, desc, col) in enumerate(next_items):
    by = Inches(1.7 + i * 1.2)
    add_rect(s, Inches(6.6), by, Inches(1.3), Inches(0.95), fill_rgb=col)
    add_text(s, priority, Inches(6.62), by + Inches(0.27), Inches(1.26), Inches(0.4),
             bold=True, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(7.9), by, Inches(5.1), Inches(0.95),
             fill_rgb=LIGHT_BG, line_rgb=col, line_pt=1)
    add_text(s, desc, Inches(8.0), by + Inches(0.1), Inches(4.9), Inches(0.75),
             size=10.5, color=DARK_TEXT)

# domanda di ricerca
add_rect(s, Inches(0.4), Inches(6.65), Inches(12.7), Inches(0.68),
         fill_rgb=NAVY)
add_text(s, "Domanda aperta: il meta-learning con Prototypical Networks riesce a superare il "
            "bottleneck ε²=0.85 adattandosi con K trial per soggetto, dove gli approcci cross-subject e subject-specific hanno fallito?",
         Inches(0.55), Inches(6.7), Inches(12.5), Inches(0.58),
         bold=True, size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_note(s, "Slide conclusiva.\nLo stato attuale è che abbiamo una documentazione completa e rigorosa del fallimento cross-subject, con 4 prove indipendenti convergenti.\nQuesto in realtà è un contributo in sé: dimostrare che nessuna architettura GNN (statica, dinamica, ipergraph) riesce a superare il problema ε²=0.85 è un risultato rilevante.\nIl meta-learning è il naturale passo successivo, motivato sia dai risultati negativi che dalla letteratura (Bomatter 2025, FAST 2025).\nTimeline: se EEG_12 funziona, la tesi ha una struttura completa: baseline → GNN → HGNN (fallimento sistematico) → meta-learning (soluzione).")

# ════════════════════════════════════════════════════════════════════════════
# SALVATAGGIO
# ════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Salvato: {OUT}")
print(f"Slide: {len(prs.slides)}")
